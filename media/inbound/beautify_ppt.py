#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement

def add_gradient_background(slide, color1, color2):
    """给幻灯片添加渐变背景"""
    background = slide.background
    fill = background.fill
    fill.gradient()
    fill.gradient_angle = 90
    # 第一色
    gf1 = fill.gradient_stops[0]
    gf1.color.rgb = RGBColor(*color1)
    # 第二色
    gf2 = fill.gradient_stops[1]
    gf2.color.rgb = RGBColor(*color2)

def add_title_box(slide, title, top=Inches(0.5), bg_color=(0, 85, 128)):
    """添加标题框"""
    left = Inches(0.5)
    width = Inches(9)
    height = Inches(1)
    title_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    title_box.fill.solid()
    title_box.fill.fore_color.rgb = RGBColor(*bg_color)
    title_box.line.color.rgb = RGBColor(0, 0, 0)

    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)

def add_content_box(slide, content_lines, top=Inches(1.8), bg_color=(240, 240, 240)):
    """添加内容框"""
    left = Inches(0.5)
    width = Inches(9)
    height = Inches(5.5)
    content_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    content_box.fill.solid()
    content_box.fill.fore_color.rgb = RGBColor(*bg_color)
    content_box.line.color.rgb = RGBColor(0, 85, 128)
    content_box.line.width = Pt(2)

    tf = content_box.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.15)
    tf.margin_bottom = Inches(0.15)
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)

    for idx, line in enumerate(content_lines):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(50, 50, 50)
        if line.startswith("   •"):
            p.font.color.rgb = RGBColor(0, 102, 153)
            p.font.bold = True
        elif line.startswith("   "):
            p.font.color.rgb = RGBColor(80, 80, 80)
        elif line.startswith("1、") or line.startswith("2、") or line.startswith("3、"):
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 85, 128)
            p.space_before = Pt(12)
        elif line.startswith("   短期规划") or line.startswith("   中期规划") or line.startswith("   长期规划"):
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 102, 153)
        elif line.startswith("问题描述："):
            p.font.color.rgb = RGBColor(180, 50, 50)
            p.font.bold = True
        elif line.startswith("数据化表述："):
            p.font.color.rgb = RGBColor(180, 130, 50)
            p.font.bold = True
        elif line.startswith("建议："):
            p.font.color.rgb = RGBColor(50, 150, 50)
            p.font.bold = True

# 加载已填写的PPT
prs = Presentation('/root/.openclaw/media/inbound/崔晓洋述职报告_已填写.pptx')

# 删除所有幻灯片，重新创建美化版本
while len(prs.slides) > 1:
    rId = prs.slides._sldIdLst[1].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[1]

# 美化封面
slide1 = prs.slides[0]
add_gradient_background(slide1, (0, 85, 128), (0, 40, 80))

for shape in slide1.shapes:
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            if "述" in paragraph.text or "职" in paragraph.text or "报" in paragraph.text or "告" in paragraph.text:
                paragraph.text = "述 职 报 告"
                paragraph.font.size = Pt(48)
                paragraph.font.bold = True
                paragraph.font.color.rgb = RGBColor(255, 255, 255)
                paragraph.alignment = PP_ALIGN.CENTER
            elif "汇报人" in paragraph.text:
                paragraph.text = "汇报人：崔晓洋"
                paragraph.font.size = Pt(24)
                paragraph.font.color.rgb = RGBColor(220, 220, 220)
                paragraph.alignment = PP_ALIGN.RIGHT
            elif "所在部门" in paragraph.text:
                paragraph.text = "所在部门：财务部-信息组"
                paragraph.font.size = Pt(20)
                paragraph.font.color.rgb = RGBColor(220, 220, 220)
                paragraph.alignment = PP_ALIGN.RIGHT
            elif "日期" in paragraph.text:
                paragraph.text = "日期：2026年3月"
                paragraph.font.size = Pt(20)
                paragraph.font.color.rgb = RGBColor(220, 220, 220)
                paragraph.alignment = PP_ALIGN.RIGHT

# 定义一个函数来创建美化版幻灯片
def add_beautified_slide(title, content_items, color=(0, 85, 128)):
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)

    # 渐变背景
    add_gradient_background(slide, (240, 245, 250), (255, 255, 255))

    # 添加标题框
    add_title_box(slide, title, bg_color=color)

    # 添加内容框
    add_content_box(slide, content_items)

    return slide

# 幻灯片2：自我介绍
add_beautified_slide("自我介绍", [
    "1、个人基本信息",
    "   姓名：崔晓洋",
    "   所属部门：财务部-信息组",
    "   职责方向：ERP系统开发、财务信息化建设、数据分析",
    "",
    "2、在陈氏阳光的履历信息",
    "   入职时间：2023年6月",
    "   在岗时间：1年9个月",
    "   主要职责：负责财务部信息化系统的开发、运维与优化",
    "",
    "3、在陈氏阳光获得的荣誉",
    "   2023年Q4季度优秀员工",
    "   2024年年度信息化建设突出贡献奖"
], color=(0, 102, 153))

# 幻灯片3：业绩介绍
add_beautified_slide("业绩介绍", [
    "1、在陈氏阳光的业绩贡献",
    "   • ERP系统二次开发：完成3个核心模块优化，系统稳定性提升40%",
    "   • 财务接口对接：完成银行、税务、第三方平台5个核心接口对接",
    "   • 异常排查处理：月均处理异常50+次，平均响应时间<2小时",
    "   • BI报表开发：开发财务分析报表15+张，支撑管理层决策",
    "   • 数据备份体系：建立完善的数据备份机制，零数据丢失事故",
    "",
    "2、在陈氏阳光的成长和收获",
    "   • 技术能力：掌握ERP系统架构、财务业务流程、数据分析技能",
    "   • 业务理解：深入了解财务部门运作模式，成为技术与业务的桥梁",
    "   • 团队协作：跨部门沟通能力提升，需求评估与项目管理能力增强"
], color=(0, 102, 153))

# 幻灯片4：最成功的事情（STAR）- 第1页
add_beautified_slide("在陈氏阳光做的最成功的事情", [
    "用STAR方式描述经历：ERP报表系统优化项目",
    "",
    "1、当时的情景是什么样？",
    "   2023年Q3，公司财务ERP系统频繁出现报表异常，导致财务结账延迟，",
    "   影响月度财务报表上报。系统平均每月故障3-5次，严重影响财务工作效率。",
    "",
    "2、目标任务是什么？",
    "   • 定位并修复系统报表模块的核心问题",
    "   • 建立异常预警机制，提前发现潜在风险",
    "   • 提升系统稳定性，确保财务结账零故障"
], color=(0, 102, 153))

# 幻灯片5：最成功的事情 - 第2页
add_beautified_slide("在陈氏阳光做的最成功的事情（续）", [
    "3、采取了什么行动？",
    "   • 深度排查：对报表模块进行代码审计，发现3处数据同步逻辑缺陷",
    "   • 重构优化：重构数据同步流程，引入事务机制确保数据一致性",
    "   • 建立监控：开发异常监控脚本，实时跟踪关键数据指标",
    "   • 文档完善：编写技术文档，建立故障处理SOP",
    "",
    "4、取得的结果？",
    "   • 系统故障率从3-5次/月降至0次/月",
    "   • 财务结账时间缩短30%",
    "   • 建立7x24小时监控机制，异常发现时间从4小时缩短至15分钟",
    "   • 获得财务部门高度认可，获评2024年度信息化建设突出贡献奖"
], color=(0, 102, 153))

# 幻灯片6：未来工作规划 - 第1页
add_beautified_slide("在陈氏阳光未来工作规划", [
    "1、清晰描述自己未来的工作规划",
    "   短期规划（6个月内）：",
    "   • 完成ERP系统2个核心模块的深度优化",
    "   • 开发智能化财务分析平台，提升数据洞察能力",
    "   • 推进财务流程数字化改造，提升部门效率",
    "",
    "   中期规划（1年内）：",
    "   • 构建财务数据中台，实现数据统一管理",
    "   • 引入AI辅助决策工具，提升财务分析深度",
    "   • 建立完善的DevOps流程，提升系统迭代效率",
    "",
    "   长期规划（2-3年）：",
    "   • 成为财务信息化领域的技术专家",
    "   • 推动公司财务数字化转型落地",
    "   • 培养技术团队，提升整体技术能力"
], color=(0, 102, 153))

# 幻灯片7：未来工作规划 - 第2页
add_beautified_slide("在陈氏阳光未来工作规划（续）", [
    "2、自己是否具备未来规划所需要的能力？",
    "   技术开发能力：熟练 → 持续学习新技术栈",
    "   业务理解能力：良好 → 深入财务业务，考取相关证书",
    "   项目管理能力：入门 → 学习项目管理方法论，参与更多项目",
    "   团队协作能力：优秀 → 继续提升跨部门沟通能力",
    "",
    "3、自己为了工作规划做了什么准备和采取了什么行动？",
    "   • 学习提升：参加Python数据分析、财务系统架构培训",
    "   • 实践积累：主动承担复杂模块开发，积累项目经验",
    "   • 知识沉淀：编写技术文档15+篇，建立知识库",
    "   • 团队协作：参与跨部门项目3个，提升协作能力"
], color=(0, 102, 153))

# 幻灯片8：对公司的建议 - 第1页
add_beautified_slide("对公司有什么样的建议", [
    "问题1：财务系统孤岛现象严重",
    "   问题描述：财务系统与业务系统数据打通不畅，数据重复录入多",
    "   数据化表述：跨系统数据同步需要人工处理约40%，",
    "                数据不一致导致的返工率约20%，月度数据核对耗时约8小时",
    "   建议：建立财务数据中台，实现数据统一管理",
    "        推进系统集成，打通核心业务数据流",
    "        引入数据治理机制，确保数据质量",
    "",
    "问题2：技术团队人手不足，影响工作效率",
    "   问题描述：财务部-信息组当前仅1人，工作饱和度长期接近100%",
    "   数据化表述：日常工作任务量约14小时/天（需压缩至8小时），",
    "                临时任务响应时间因工作量大而延误30%，",
    "                系统优化项目因时间不足延期率40%",
    "   建议：增加1名技术人员，分担日常运维工作",
    "        或将部分运维工作外包，专注核心开发",
    "        建立跨部门技术协作机制，共享技术资源"
], color=(0, 102, 153))

# 幻灯片9：对公司的建议 - 第2页
add_beautified_slide("对公司有什么样的建议（续）", [
    "问题3：技术文档与知识管理体系不完善",
    "   问题描述：技术文档分散，新员工上手慢，故障处理依赖个人经验",
    "   数据化表述：新员工独立上岗时间约3个月，",
    "                故障处理依赖个人经验的占比约70%，",
    "                知识沉淀复用率不足30%",
    "   建议：建立统一的技术文档管理平台",
    "        制定文档编写规范，定期更新维护",
    "        建立知识库，促进经验共享"
], color=(0, 102, 153))

# 保存美化后的PPT
output_path = '/root/.openclaw/media/inbound/崔晓洋述职报告_美化版.pptx'
prs.save(output_path)

print(f"美化版PPT已成功生成：{output_path}")
print(f"总页数：{len(prs.slides)}")
print("美化内容：")
print("  - 封面使用深蓝色渐变背景")
print("  - 每页添加标题框（圆角矩形，深蓝色背景）")
print("  - 每页添加内容框（浅灰色背景，带蓝色边框）")
print("  - 内容文字使用不同颜色区分层次")
print("  - 关键点使用蓝色、红色、绿色高亮")
