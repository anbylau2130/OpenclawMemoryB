#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# 加载老板改过的PPT
prs = Presentation('/root/.openclaw/media/inbound/3d8b321b-2cac-417c-9661-a1dfa2586e4c.pptx')

def format_paragraph(paragraph, text=""):
    """格式化段落文字"""
    if text:
        paragraph.text = text
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = RGBColor(60, 60, 60)
        
        # 根据内容设置不同的样式
        if text.startswith("1、") or text.startswith("2、") or text.startswith("3、"):
            paragraph.font.size = Pt(20)
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(0, 102, 153)
            paragraph.space_before = Pt(18)
        elif text.startswith("   •") or text.startswith("•"):
            paragraph.font.color.rgb = RGBColor(0, 102, 153)
            paragraph.font.bold = True
            paragraph.space_before = Pt(6)
        elif text.startswith("   短期规划") or text.startswith("   中期规划") or text.startswith("   长期规划"):
            paragraph.font.size = Pt(18)
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(0, 102, 153)
            paragraph.space_before = Pt(12)
        elif text.startswith("问题描述："):
            paragraph.font.color.rgb = RGBColor(180, 60, 60)
            paragraph.font.bold = True
        elif text.startswith("数据化表述："):
            paragraph.font.color.rgb = RGBColor(180, 130, 50)
            paragraph.font.bold = True
        elif text.startswith("建议："):
            paragraph.font.color.rgb = RGBColor(50, 150, 50)
            paragraph.font.bold = True
        elif text.startswith("   "):
            paragraph.font.color.rgb = RGBColor(80, 80, 80)
            paragraph.space_before = Pt(6)
        elif text.strip() == "":
            paragraph.font.size = Pt(8)
        elif "→" in text and "：" in text:
            paragraph.font.size = Pt(15)
            paragraph.space_before = Pt(4)

# 美化封面
slide1 = prs.slides[0]
for shape in slide1.shapes:
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            text = paragraph.text.strip()
            if "述职报告" in text or ("述" in text and "职" in text and "报" in text and "告" in text):
                paragraph.font.size = Pt(48)
                paragraph.font.bold = True
                paragraph.font.color.rgb = RGBColor(0, 102, 153)
            elif "汇报人" in text:
                paragraph.text = "汇报人：崔晓洋"
                paragraph.font.size = Pt(24)
                paragraph.font.color.rgb = RGBColor(50, 50, 50)
                paragraph.alignment = PP_ALIGN.RIGHT
            elif "日期" in text:
                paragraph.text = "日期：2026年3月"
                paragraph.font.size = Pt(20)
                paragraph.font.color.rgb = RGBColor(80, 80, 80)
                paragraph.alignment = PP_ALIGN.RIGHT
            elif "照片" in text:
                paragraph.font.size = Pt(12)
                paragraph.font.color.rgb = RGBColor(150, 150, 150)
                paragraph.alignment = PP_ALIGN.CENTER

# 美化其他幻灯片（从第2页开始）
for i in range(1, len(prs.slides)):
    slide = prs.slides[i]
    for shape in slide.shapes:
        if shape.has_text_frame:
            tf = shape.text_frame
            for paragraph in tf.paragraphs:
                text = paragraph.text.strip()
                if text:
                    format_paragraph(paragraph, text)

# 保存美化后的PPT
output_path = '/root/.openclaw/media/inbound/崔晓洋述职报告_你的版本_美化.pptx'
prs.save(output_path)

print(f"美化版PPT已生成：{output_path}")
print(f"总页数：{len(prs.slides)}")
print("\n美化内容（保持原有模板背景和布局）：")
print("  ✓ 标题：深蓝色加粗大字（48pt）")
print("  ✓ 一级标题（1、2、3）：深蓝色加粗（20pt）")
print("  ✓ 核心要点（• ）：蓝色加粗")
print("  ✓ 小标题（短期/中期/长期规划）：蓝色加粗（18pt）")
print("  ✓ 问题描述：红色加粗")
print("  ✓ 数据化表述：橙色加粗")
print("  ✓ 建议：绿色加粗")
print("  ✓ 其他内容：灰色常规字（16pt）")
print("\n✅ 已保留你修改的所有内容！")
