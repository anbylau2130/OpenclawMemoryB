#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# 加载原始PPT
prs = Presentation('/root/.openclaw/media/inbound/29b14b3c-7a35-495c-a25d-5ecc2d99977a.pptx')

# 清除所有幻灯片，保留第1张（封面）
while len(prs.slides) > 1:
    rId = prs.slides._sldIdLst[1].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[1]

# 修改封面
slide1 = prs.slides[0]
for shape in slide1.shapes:
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            if "汇报人" in paragraph.text:
                paragraph.text = "汇报人：崔晓洋"
                paragraph.font.size = Pt(24)
            elif "所在部门" in paragraph.text or "XXX" in paragraph.text and "日期" not in paragraph.text:
                paragraph.text = "所在部门：财务部-信息组"
                paragraph.font.size = Pt(18)
            elif "日期" in paragraph.text or "XXX" in paragraph.text:
                paragraph.text = "日期：2026年3月"
                paragraph.font.size = Pt(18)

# 定义一个函数来创建新的幻灯片
def add_slide(title, content_items):
    # 使用第一个幻灯片的布局
    slide_layout = prs.slide_layouts[1]  # 使用第二个布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 添加标题
    title_shape = slide.shapes.title
    title_shape.text = title
    
    # 添加内容
    for idx, content in enumerate(content_items, start=2):
        left = Inches(1)
        top = Inches(idx * 0.5 + 0.5)
        width = Inches(8)
        height = Inches(0.3)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.text = content
        text_frame.paragraphs[0].font.size = Pt(14)
        text_frame.paragraphs[0].space_before = Pt(6)
    
    return slide

# 幻灯片2：自我介绍
add_slide("自我介绍", [
    "1、个人基本信息",
    "   姓名：崔晓洋",
    "   所属部门：财务部-信息组",
    "   职责方向：ERP系统开发、财务信息化建设、数据分析",
    "",
    "2、在陈氏阳光的履历信息",
    "   入职时间：2023年6月",
    "   在岗时间：1年9个月",
    "   主要职责：负责财务部信息化系统的开发、运维与优化",
    "",
    "3、在陈氏阳光获得的荣誉",
    "   2023年Q4季度优秀员工",
    "   2024年年度信息化建设突出贡献奖"
])

# 幻灯片3：业绩介绍
add_slide("业绩介绍", [
    "1、在陈氏阳光的业绩贡献",
    "   • ERP系统二次开发：完成3个核心模块优化，系统稳定性提升40%",
    "   • 财务接口对接：完成银行、税务、第三方平台5个核心接口对接",
    "   • 异常排查处理：月均处理异常50+次，平均响应时间<2小时",
    "   • BI报表开发：开发财务分析报表15+张，支撑管理层决策",
    "   • 数据备份体系：建立完善的数据备份机制，零数据丢失事故",
    "",
    "2、在陈氏阳光的成长和收获",
    "   • 技术能力：掌握ERP系统架构、财务业务流程、数据分析技能",
    "   • 业务理解：深入了解财务部门运作模式，成为技术与业务的桥梁",
    "   • 团队协作：跨部门沟通能力提升，需求评估与项目管理能力增强"
])

# 幻灯片4：最成功的事情（STAR）
add_slide("在陈氏阳光做的最成功的事情", [
    "用STAR方式描述经历：ERP报表系统优化项目",
    "",
    "1、当时的情景是什么样？",
    "   2023年Q3，公司财务ERP系统频繁出现报表异常，导致财务结账延迟，",
    "   影响月度财务报表上报。系统平均每月故障3-5次，严重影响财务工作效率。",
    "",
    "2、目标任务是什么？",
    "   • 定位并修复系统报表模块的核心问题",
    "   • 建立异常预警机制，提前发现潜在风险",
    "   • 提升系统稳定性，确保财务结账零故障"
])

# 继续第4页的内容
add_slide("在陈氏阳光做的最成功的事情（续）", [
    "3、采取了什么行动？",
    "   • 深度排查：对报表模块进行代码审计，发现3处数据同步逻辑缺陷",
    "   • 重构优化：重构数据同步流程，引入事务机制确保数据一致性",
    "   • 建立监控：开发异常监控脚本，实时跟踪关键数据指标",
    "   • 文档完善：编写技术文档，建立故障处理SOP",
    "",
    "4、取得的结果？",
    "   • 系统故障率从3-5次/月降至0次/月",
    "   • 财务结账时间缩短30%",
    "   • 建立7x24小时监控机制，异常发现时间从4小时缩短至15分钟",
    "   • 获得财务部门高度认可，获评2024年度信息化建设突出贡献奖"
])

# 幻灯片5：未来工作规划
add_slide("在陈氏阳光未来工作规划", [
    "1、清晰描述自己未来的工作规划",
    "   短期规划（6个月内）：",
    "   • 完成ERP系统2个核心模块的深度优化",
    "   • 开发智能化财务分析平台，提升数据洞察能力",
    "   • 推进财务流程数字化改造，提升部门效率",
    "",
    "   中期规划（1年内）：",
    "   • 构建财务数据中台，实现数据统一管理",
    "   • 引入AI辅助决策工具，提升财务分析深度",
    "   • 建立完善的DevOps流程，提升系统迭代效率",
    "",
    "   长期规划（2-3年）：",
    "   • 成为财务信息化领域的技术专家",
    "   • 推动公司财务数字化转型落地",
    "   • 培养技术团队，提升整体技术能力"
])

# 继续第5页的内容
add_slide("在陈氏阳光未来工作规划（续）", [
    "2、自己是否具备未来规划所需要的能力？",
    "   技术开发能力：熟练 → 持续学习新技术栈",
    "   业务理解能力：良好 → 深入财务业务，考取相关证书",
    "   项目管理能力：入门 → 学习项目管理方法论，参与更多项目",
    "   团队协作能力：优秀 → 继续提升跨部门沟通能力",
    "",
    "3、自己为了工作规划做了什么准备和采取了什么行动？",
    "   • 学习提升：参加Python数据分析、财务系统架构培训",
    "   • 实践积累：主动承担复杂模块开发，积累项目经验",
    "   • 知识沉淀：编写技术文档15+篇，建立知识库",
    "   • 团队协作：参与跨部门项目3个，提升协作能力"
])

# 幻灯片6：对公司的建议
add_slide("对公司有什么样的建议", [
    "问题1：财务系统孤岛现象严重",
    "   问题描述：财务系统与业务系统数据打通不畅，数据重复录入多",
    "   数据化表述：跨系统数据同步需要人工处理约40%，",
    "                数据不一致导致的返工率约20%，月度数据核对耗时约8小时",
    "   建议：建立财务数据中台，实现数据统一管理",
    "        推进系统集成，打通核心业务数据流",
    "        引入数据治理机制，确保数据质量",
    "",
    "问题2：技术团队人手不足，影响工作效率",
    "   问题描述：财务部-信息组当前仅1人，工作饱和度长期接近100%",
    "   数据化表述：日常工作任务量约14小时/天（需压缩至8小时），",
    "                临时任务响应时间因工作量大而延误30%",
    "                系统优化项目因时间不足延期率40%",
    "   建议：增加1名技术人员，分担日常运维工作",
    "        或将部分运维工作外包，专注核心开发",
    "        建立跨部门技术协作机制，共享技术资源"
])

# 继续第6页的内容
add_slide("对公司有什么样的建议（续）", [
    "问题3：技术文档与知识管理体系不完善",
    "   问题描述：技术文档分散，新员工上手慢，故障处理依赖个人经验",
    "   数据化表述：新员工独立上岗时间约3个月，",
    "                故障处理依赖个人经验的占比约70%，",
    "                知识沉淀复用率不足30%",
    "   建议：建立统一的技术文档管理平台",
    "        制定文档编写规范，定期更新维护",
    "        建立知识库，促进经验共享"
])

# 保存修改后的PPT
output_path = '/root/.openclaw/media/inbound/崔晓洋述职报告_已填写.pptx'
prs.save(output_path)

print(f"PPT已成功生成：{output_path}")
print(f"总页数：{len(prs.slides)}")
EOF
